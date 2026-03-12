#!/usr/bin/env python3
"""
Generate HPE-branded PowerPoint presentation for HPE AI Autopilot demo.
Uses standard HPE brand colors: HPE Green (#01A982), dark backgrounds, white text.

Refactored: each slide is a function, layout constants extracted, slide numbers added.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =============================================================================
# HPE Brand Colors
# =============================================================================
HPE_GREEN = RGBColor(0x01, 0xA9, 0x82)
HPE_DARK = RGBColor(0x1A, 0x1A, 0x2E)
HPE_DARKER = RGBColor(0x12, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)

# =============================================================================
# Layout Constants (extracted from repeated magic numbers)
# =============================================================================
LEFT_MARGIN = Inches(0.8)
RIGHT_COL = Inches(7.2)
CARD_BG = RGBColor(0x24, 0x27, 0x36)
CARD_BORDER = RGBColor(0x2E, 0x32, 0x41)
ROW_ALT_BG = RGBColor(0x1E, 0x20, 0x2E)
CONTENT_WIDTH = Inches(11.7)
ACCENT_WIDTH = Inches(0.15)

TOTAL_SLIDES = 11

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "..", "screenshots")


# =============================================================================
# Helper Functions
# =============================================================================

def add_hpe_background(slide, color=HPE_DARK):
    """Fill slide background with HPE dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_hpe_header_bar(slide, W):
    """Add the HPE green accent bar at the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        W, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HPE_GREEN
    shape.line.fill.background()


def add_hpe_footer(slide, W, H, slide_num=None, show_confidential=True):
    """Add HPE footer with logo text, confidential notice, and slide number."""
    # Bottom bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), H - Inches(0.5),
        W, Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = HPE_DARKER
    bar.line.fill.background()

    # HPE Logo text (left)
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  Hewlett Packard Enterprise"
    run.font.size = Pt(10)
    run.font.color.rgb = HPE_GREEN
    run.font.bold = True

    # Slide number (center)
    if slide_num is not None:
        num_box = slide.shapes.add_textbox(
            W / 2 - Inches(0.5), H - Inches(0.45),
            Inches(1.0), Inches(0.4)
        )
        tf2 = num_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = f"{slide_num} / {TOTAL_SLIDES}"
        r2.font.size = Pt(9)
        r2.font.color.rgb = MEDIUM_GRAY

    if show_confidential:
        # Confidential text (right side)
        conf = slide.shapes.add_textbox(
            W - Inches(3), H - Inches(0.45),
            Inches(2.8), Inches(0.4)
        )
        tf3 = conf.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        r3.text = "Internal Use Only"
        r3.font.size = Pt(9)
        r3.font.color.rgb = MEDIUM_GRAY
        r3.font.italic = True


def add_green_accent_shape(slide, left, top, width, height):
    """Add a subtle green accent rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HPE_GREEN
    shape.line.fill.background()
    return shape


def add_slide_boilerplate(prs, W, H, slide_num, bg_color=HPE_DARK):
    """Common slide setup: background, header bar, green accent, footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_hpe_background(slide, bg_color)
    add_hpe_header_bar(slide, W)
    add_green_accent_shape(slide, Inches(0), Inches(0.06), ACCENT_WIDTH, H - Inches(0.56))
    return slide


def add_slide_title(slide, text, size=Pt(36)):
    """Add a standard slide title."""
    title_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = WHITE
    run.font.bold = True


def add_section_header(slide, text, left, top, width=Inches(5.5)):
    """Add a green uppercase section header with letter spacing."""
    hdr = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True
    r.font.letter_spacing = Pt(2)


def add_vertical_divider(slide, left, top, height=Inches(4.8)):
    """Add a subtle vertical divider line."""
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top,
        Inches(0.03), height
    )
    div.fill.solid()
    div.fill.fore_color.rgb = CARD_BORDER
    div.line.fill.background()


def add_card(slide, left, top, width, height, border_color=None):
    """Add a dark card shape with optional custom border color."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_color or CARD_BORDER
    card.line.width = Pt(1.5) if border_color else Pt(1)
    return card


# =============================================================================
# SLIDE 1 — Title Slide
# =============================================================================
def create_slide_1_title(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 1)

    # HPE Logo text
    logo_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), Inches(4), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Hewlett Packard Enterprise"
    run.font.size = Pt(14)
    run.font.color.rgb = HPE_GREEN
    run.font.bold = True

    # Main title
    title_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.0), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "HPE AI Autopilot"
    run.font.size = Pt(48)
    run.font.color.rgb = WHITE
    run.font.bold = True

    # Subtitle
    sub_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.5), Inches(10), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Intelligent IT Operations — From Alert to Resolution in Minutes"
    run.font.size = Pt(22)
    run.font.color.rgb = LIGHT_GRAY

    # Separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN, Inches(4.6),
        Inches(3), Inches(0.04)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = HPE_GREEN
    sep.line.fill.background()

    # Presenter info
    info_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.9), Inches(6), Inches(1.2))
    tf = info_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "AI Innovation Demo  |  March 2026"
    r1.font.size = Pt(16)
    r1.font.color.rgb = MEDIUM_GRAY

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Vishnu Dutt Kidambi"
    r2.font.size = Pt(14)
    r2.font.color.rgb = MEDIUM_GRAY

    # Right-side — Key Capabilities (labeled explicitly, with count note)
    cap_hdr = slide.shapes.add_textbox(Inches(9.5), Inches(1.4), Inches(3.5), Inches(0.4))
    tf = cap_hdr.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "KEY CAPABILITIES"
    r.font.size = Pt(11)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True
    r.font.letter_spacing = Pt(2)

    tools_box = slide.shapes.add_textbox(Inches(9.5), Inches(1.8), Inches(3.5), Inches(4.5))
    tf = tools_box.text_frame
    tf.word_wrap = True

    capabilities = [
        ("\U0001f534", "Alert Intelligence"),
        ("\U0001f50d", "Deep Investigation"),
        ("\U0001f50c", "Network Correlation"),
        ("\U0001f4a5", "Blast Radius Analysis"),
        ("\U0001f4c8", "Capacity Forecasting"),
        ("\U0001f4d6", "Runbook Search"),
        ("\U0001f6e0\ufe0f", "Guided Remediation"),
    ]

    for icon, label in capabilities:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)

        ri = p.add_run()
        ri.text = f"{icon}  "
        ri.font.size = Pt(16)

        rl = p.add_run()
        rl.text = label
        rl.font.size = Pt(15)
        rl.font.color.rgb = LIGHT_GRAY

    # Note about tool count
    note_p = tf.add_paragraph()
    note_p.space_before = Pt(10)
    note_r = note_p.add_run()
    note_r.text = "(7 of 11 tools shown)"
    note_r.font.size = Pt(10)
    note_r.font.color.rgb = MEDIUM_GRAY
    note_r.font.italic = True

    add_hpe_footer(slide, W, H, slide_num=1)


# =============================================================================
# SLIDE 2 — The Problem + Solution
# =============================================================================
def create_slide_2_problem_solution(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 2)
    add_slide_title(slide, "Why AI Autopilot?")

    # Left column — The Problem
    add_section_header(slide, "THE CHALLENGE", LEFT_MARGIN, Inches(1.6))

    problems = [
        ("Alert Fatigue", "Operators sift through hundreds of alerts manually to find what matters"),
        ("Slow Triage", "Root cause identification takes 30-60 min across multiple consoles"),
        ("Siloed Data", "Server metrics, network telemetry, and incidents live in separate tools"),
        ("Manual Runbooks", "Engineers search wikis and PDFs for resolution steps under pressure"),
    ]

    prob_y = Inches(2.2)
    for title, desc in problems:
        box = slide.shapes.add_textbox(LEFT_MARGIN, prob_y, Inches(5.5), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"\u2717  {title}"
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)  # red
        r.font.bold = True

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = f"     {desc}"
        r2.font.size = Pt(13)
        r2.font.color.rgb = LIGHT_GRAY

        prob_y += Inches(1.05)

    # Vertical divider
    add_vertical_divider(slide, Inches(6.6), Inches(1.6))

    # Right column — The Solution
    add_section_header(slide, "AI AUTOPILOT SOLUTION", RIGHT_COL, Inches(1.6))

    solutions = [
        ("Natural Language Ops", "Ask questions in plain English \u2014 the AI calls the right tools automatically"),
        ("End-to-End Investigation", "Alerts \u2192 Server \u2192 Network \u2192 Blast Radius in one conversation"),
        ("Unified Correlation", "Cross-layer analysis: OpsRamp metrics + Juniper network telemetry"),
        ("Instant Runbooks & Remediation", "AI retrieves procedures and generates CLI fix commands on demand"),
    ]

    sol_y = Inches(2.2)
    for title, desc in solutions:
        box = slide.shapes.add_textbox(RIGHT_COL, sol_y, Inches(5.5), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"\u2713  {title}"
        r.font.size = Pt(16)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = f"     {desc}"
        r2.font.size = Pt(13)
        r2.font.color.rgb = LIGHT_GRAY

        sol_y += Inches(1.05)

    # Bottom callout
    callout = add_card(slide, LEFT_MARGIN, Inches(6.3), CONTENT_WIDTH, Inches(0.6), border_color=HPE_GREEN)

    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u26a1  MTTR Reduction Target: 30-60 minutes \u2192 under 5 minutes with AI-assisted triage and guided remediation"
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.bold = True

    add_hpe_footer(slide, W, H, slide_num=2)


# =============================================================================
# SLIDE 3 — Architecture & Data Sources
# =============================================================================
def create_slide_3_architecture(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 3)
    add_slide_title(slide, "Architecture & Data Sources")

    # Architecture section (left)
    add_section_header(slide, "ARCHITECTURE", LEFT_MARGIN, Inches(1.5))

    # Architecture diagram as text boxes — FIXED: 11 tools listed correctly
    arch_items = [
        ("User", "Web UI  /  MCP Client (VS Code, Claude Desktop)  /  Terminal REPL", LEFT_MARGIN, Inches(2.0)),
        ("LLM Agent", "Ollama (llama3.1) \u2014 Tool-calling orchestration loop with streaming", LEFT_MARGIN, Inches(2.8)),
        ("11 Tools", "Alerts \u2022 Resources \u2022 Resource Details \u2022 Incidents \u2022 Investigation \u2022 Summary\nCapacity \u2022 KB Search \u2022 Network Correlation \u2022 Blast Radius \u2022 Remediation", LEFT_MARGIN, Inches(3.6)),
        ("Data Layer", "OpsRamp API (mocked)  |  Juniper Mist API (simulated)  |  PDF Runbooks (RAG)", LEFT_MARGIN, Inches(4.4)),
    ]

    for label, desc, x, y in arch_items:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(5.5), Inches(0.65)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        lbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.05), Inches(1.4), Inches(0.3))
        tf = lbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(12)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True

        dbox = slide.shapes.add_textbox(x + Inches(1.6), y + Inches(0.05), Inches(3.7), Inches(0.55))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = desc
        r.font.size = Pt(11)
        r.font.color.rgb = LIGHT_GRAY

        if label != "Data Layer":
            arrow = slide.shapes.add_textbox(x + Inches(2.5), y + Inches(0.6), Inches(0.5), Inches(0.3))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "\u2193"
            r.font.size = Pt(16)
            r.font.color.rgb = MEDIUM_GRAY

    # Divider
    add_vertical_divider(slide, Inches(6.8), Inches(1.5))

    # Data Sources section (right)
    add_section_header(slide, "DATA SOURCES", RIGHT_COL, Inches(1.5))

    # OpsRamp API card
    add_card(slide, RIGHT_COL, Inches(2.1), Inches(5.5), Inches(1.5))

    ops_box = slide.shapes.add_textbox(Inches(7.4), Inches(2.15), Inches(5.2), Inches(1.4))
    tf = ops_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "OpsRamp APIs \u2014 Mocked"
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE
    r.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "Based on the public OpsRamp V3 developer API:"
    r2.font.size = Pt(12)
    r2.font.color.rgb = LIGHT_GRAY

    p3 = tf.add_paragraph()
    p3.space_before = Pt(2)
    r3 = p3.add_run()
    r3.text = "https://develop.opsramp.com/v3"
    r3.font.size = Pt(12)
    r3.font.color.rgb = HPE_GREEN
    r3.font.bold = True

    p4 = tf.add_paragraph()
    p4.space_before = Pt(4)
    r4 = p4.add_run()
    r4.text = "Alerts, Resources, Incidents, and Metric History are mocked in-memory using realistic schemas matching the V3 API response format."
    r4.font.size = Pt(11)
    r4.font.color.rgb = MEDIUM_GRAY

    # Juniper API card
    add_card(slide, RIGHT_COL, Inches(3.8), Inches(5.5), Inches(1.5))

    jun_box = slide.shapes.add_textbox(Inches(7.4), Inches(3.85), Inches(5.2), Inches(1.4))
    tf = jun_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Juniper Mist APIs \u2014 Simulated"
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE
    r.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = "Modeled after the Juniper Mist Switch Stats API:"
    r2.font.size = Pt(11)
    r2.font.color.rgb = LIGHT_GRAY

    p3 = tf.add_paragraph()
    p3.space_before = Pt(1)
    r3 = p3.add_run()
    r3.text = "GET /api/v1/sites/{site_id}/stats/devices?type=switch"
    r3.font.size = Pt(11)
    r3.font.color.rgb = HPE_GREEN
    r3.font.bold = True

    p4 = tf.add_paragraph()
    p4.space_before = Pt(1)
    r4 = p4.add_run()
    r4.text = "Ref: Juniper Mist API (community ref: doc.mist-lab.fr)"
    r4.font.size = Pt(10)
    r4.font.color.rgb = MEDIUM_GRAY

    p5 = tf.add_paragraph()
    p5.space_before = Pt(3)
    r5 = p5.add_run()
    r5.text = "Covers: EX switch stats, port errors, packet loss, CRC errors, link flaps, latency/jitter, and dependency graph for blast radius."
    r5.font.size = Pt(10)
    r5.font.color.rgb = MEDIUM_GRAY

    # PDF Runbooks note
    pdf_note = slide.shapes.add_textbox(RIGHT_COL, Inches(5.5), Inches(5.5), Inches(0.6))
    tf = pdf_note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\U0001f4d6  PDF Runbooks \u2014 Real PDFs ingested via RAG pipeline (chunking \u2192 embedding \u2192 vector search)"
    r.font.size = Pt(12)
    r.font.color.rgb = LIGHT_GRAY

    add_hpe_footer(slide, W, H, slide_num=3)


# =============================================================================
# SLIDE 4 — OpsRamp ↔ Juniper: How Data Connects
# =============================================================================
def create_slide_4_data_connection(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 4)
    add_slide_title(slide, "OpsRamp \u2194 Juniper: How Data Connects")

    # --- Left column: IP as the Join Key ---
    add_section_header(slide, "IP ADDRESS = THE JOIN KEY", LEFT_MARGIN, Inches(1.5))

    # Visual diagram box
    add_card(slide, LEFT_MARGIN, Inches(2.0), Inches(5.5), Inches(2.2))

    diag_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.1), Inches(2.0))
    tf = diag_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "OpsRamp Resource"
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = '  Name: "app-server-prod-02"    IP: 10.0.1.103'
    r2.font.size = Pt(11)
    r2.font.color.rgb = LIGHT_GRAY

    p3 = tf.add_paragraph()
    p3.space_before = Pt(10)
    r3 = p3.add_run()
    r3.text = "         \u2500\u2500\u2500\u2500 IP 10.0.1.103 \u2500\u2500\u2500\u2500\u25b6"
    r3.font.size = Pt(13)
    r3.font.color.rgb = HPE_GREEN
    r3.font.bold = True

    p4 = tf.add_paragraph()
    p4.space_before = Pt(10)
    r4 = p4.add_run()
    r4.text = "Juniper Switch Port"
    r4.font.size = Pt(13)
    r4.font.color.rgb = WHITE
    r4.font.bold = True

    p5 = tf.add_paragraph()
    p5.space_before = Pt(2)
    r5 = p5.add_run()
    r5.text = '  Switch: sw-dc-east-01  Port: ge-0/0/3  IP: 10.0.1.103'
    r5.font.size = Pt(11)
    r5.font.color.rgb = LIGHT_GRAY

    p6 = tf.add_paragraph()
    p6.space_before = Pt(10)
    r6 = p6.add_run()
    r6.text = "The agent looks up a server\u2019s IP in OpsRamp, finds the matching\nswitch port in Juniper, and cross-correlates metrics from both."
    r6.font.size = Pt(11)
    r6.font.color.rgb = MEDIUM_GRAY

    # --- 1:1 vs Multi mapping ---
    add_section_header(slide, "PORT MAPPING SCENARIOS", LEFT_MARGIN, Inches(4.4))

    mapping_rows = [
        ("1 Server \u2192 1 Port", "Standard. Each physical server has one NIC\ncabled to one switch port. Direct correlation.", HPE_GREEN),
        ("K8s Pod \u2192 Node IP", "Pods have overlay IPs (10.244.x.x) invisible to switch.\nCorrelation uses Node IP \u2192 switch port.", HPE_GREEN),
        ("N VMs \u2192 1 Port", "Hypervisor host on one port. Many VMs share it.\nCorrelation is host-level, not per-VM.", RGBColor(0xFF, 0xC1, 0x07)),
    ]

    map_y = Inches(4.85)
    for label, desc, color in mapping_rows:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            LEFT_MARGIN, map_y, Inches(5.5), Inches(0.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        lbox = slide.shapes.add_textbox(Inches(1.0), map_y + Inches(0.05), Inches(1.9), Inches(0.45))
        tf = lbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(12)
        r.font.color.rgb = color
        r.font.bold = True

        dbox = slide.shapes.add_textbox(Inches(3.0), map_y + Inches(0.02), Inches(3.1), Inches(0.5))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = desc
        r.font.size = Pt(10)
        r.font.color.rgb = LIGHT_GRAY

        map_y += Inches(0.58)

    # --- Right column: Where Juniper correlation works ---
    add_vertical_divider(slide, Inches(6.6), Inches(1.5))

    add_section_header(slide, "WHEN OPSRAMP \u2194 JUNIPER MAPPING WORKS", RIGHT_COL, Inches(1.5))

    # On-Prem card
    add_card(slide, RIGHT_COL, Inches(2.0), Inches(5.5), Inches(2.2), border_color=HPE_GREEN)

    works_box = slide.shapes.add_textbox(Inches(7.4), Inches(2.05), Inches(5.2), Inches(2.1))
    tf = works_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\u2713  Two Prerequisites for Correlation"
    r.font.size = Pt(14)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True

    prereq_items = [
        ("1. Juniper Switches in the DC", "Physical switch fabric connecting servers \u2014 provides port stats, packet loss, link flaps, CRC errors"),
        ("2. OpsRamp Collector Deployed", "Collector gathers server-side metrics (CPU, memory, disk, alerts, incidents) from the same infra"),
    ]
    for title, desc in prereq_items:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = f"  {title}"
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE
        r.font.bold = True
        p = tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = f"     {desc}"
        r.font.size = Pt(10)
        r.font.color.rgb = LIGHT_GRAY

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "  \u2192 When both are present, OpsRamp server stats can be"
    r.font.size = Pt(11)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True
    p = tf.add_paragraph()
    p.space_before = Pt(1)
    r = p.add_run()
    r.text = "     correlated with Juniper switch stats via IP \u2192 port mapping"
    r.font.size = Pt(11)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True

    # Example card — HPE GreenLake
    add_card(slide, RIGHT_COL, Inches(4.4), Inches(5.5), Inches(1.8))

    example_box = slide.shapes.add_textbox(Inches(7.4), Inches(4.45), Inches(5.2), Inches(1.7))
    tf = example_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Example: HPE GreenLake DC"
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True

    example_items = [
        "HPE-managed datacenter with Juniper EX switches",
        "OpsRamp collector monitors K8s nodes, VMs, and apps",
        "Node IP (e.g. 10.0.4.24) maps to switch port ge-0/0/5",
        "AI correlates: high HTTP latency \u2192 port packet loss \u2192 link flaps",
    ]
    for item in example_items:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"  \u2022  {item}"
        r.font.size = Pt(11)
        r.font.color.rgb = LIGHT_GRAY

    add_hpe_footer(slide, W, H, slide_num=4)


# =============================================================================
# SLIDE 5 — End-to-End Flow: "Why is GreenLake Portal Slow?"
# =============================================================================
def create_slide_5_e2e_flow(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 5)
    add_slide_title(slide, 'End-to-End: "Why is GreenLake Portal Slow?"', size=Pt(34))

    # User prompt callout
    prompt_card = add_card(slide, LEFT_MARGIN, Inches(1.35), CONTENT_WIDTH, Inches(0.5), border_color=HPE_GREEN)

    prompt_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.38), Inches(11.3), Inches(0.45))
    tf = prompt_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '\U0001f4ac  User types: "Why is greenlake portal slow?"'
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Flow steps — 6-step investigation chain
    flow_steps = [
        ("1", "search_alerts",
         "Finds 2 critical alerts on k8s-node-04 (HPE GreenLake):\ncontainer restart loop (P0) + HTTP response 6.2s (P1)",
         "OpsRamp"),
        ("2", "investigate_resource",
         "Deep-dives k8s-node-04: CPU 60%, memory 65%, disk 45%,\n2 active alerts, 1 open incident \u2014 server metrics look OK",
         "OpsRamp"),
        ("3", "correlate_network",
         "Resolves greenlake-portal \u2192 k8s-node-04 \u2192 IP 10.0.4.24\nFinds sw-dc-east-04:ge-0/0/5: 8.3% loss, 156K errors, 100Mbps half-duplex, link flaps",
         "Juniper"),
        ("4", "blast_radius",
         "BFS traversal: k8s-node-04 \u2192 9 apps including greenlake-portal,\naruba-central, dscc-console \u2192 5,200 users affected, $75K+/hr revenue at risk",
         "Dependency\nGraph"),
        ("5", "search_knowledge_base",
         "RAG search finds runbook: 'Switch Port Troubleshooting' \u2014\nJunos CLI commands for link flap diagnosis & cable replacement",
         "PDF / RAG"),
        ("6", "get_remediation_plan",
         "Generates 8-step fix: verify port \u2192 clear counters \u2192 check optics\n\u2192 force duplex \u2192 bounce port \u2192 replace cable \u2014 with exact Junos CLI",
         "Juniper"),
    ]

    step_y = Inches(2.0)
    for num, tool, desc, source in flow_steps:
        # Step card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            LEFT_MARGIN, step_y, CONTENT_WIDTH, Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), step_y + Inches(0.13), Inches(0.44), Inches(0.44)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = HPE_GREEN
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.size = Pt(16)
        cr.font.color.rgb = WHITE
        cr.font.bold = True

        # Tool name
        tbox = slide.shapes.add_textbox(Inches(1.6), step_y + Inches(0.05), Inches(2.3), Inches(0.3))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tool
        r.font.size = Pt(13)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True

        # Description
        dbox = slide.shapes.add_textbox(Inches(3.9), step_y + Inches(0.03), Inches(7.0), Inches(0.65))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = desc
        r.font.size = Pt(10)
        r.font.color.rgb = LIGHT_GRAY

        # Source badge
        sbox = slide.shapes.add_textbox(Inches(11.0), step_y + Inches(0.15), Inches(1.3), Inches(0.4))
        tf = sbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = source
        r.font.size = Pt(9)
        r.font.color.rgb = MEDIUM_GRAY
        r.font.italic = True

        # Arrow between steps
        if num != "6":
            arrow = slide.shapes.add_textbox(Inches(1.1), step_y + Inches(0.65), Inches(0.3), Inches(0.18))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "\u2193"
            r.font.size = Pt(12)
            r.font.color.rgb = HPE_GREEN

        step_y += Inches(0.73)

    # Bottom result callout — ENHANCED with time element
    result_card = add_card(slide, LEFT_MARGIN, Inches(6.4), CONTENT_WIDTH, Inches(0.5), border_color=HPE_GREEN)

    tf = result_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u26a1  Result: Root cause = link flaps + duplex mismatch on ge-0/0/5, business impact = 5,200 users across 9 apps, and 8-step CLI fix \u2014 all in one conversation, under 2 minutes"
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True

    add_hpe_footer(slide, W, H, slide_num=5)


# =============================================================================
# SLIDE 6 — 11 Tools Overview (Part 1: Tools 1–6)
# =============================================================================
def create_slide_6_tools_part1(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 6)
    add_slide_title(slide, "11 AI Tools \u2014 Data & Investigation")

    sub6 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.15), Inches(11), Inches(0.4))
    tf = sub6.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "The LLM autonomously picks which tools to call based on the user\u2019s question"
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_GRAY

    # Tool table — Part 1
    tools_part1 = [
        ("#", "Tool", "What It Does", "Data Source"),
        ("1", "search_alerts", "Find alerts by severity, priority, resource, or keyword", "OpsRamp Alerts API (mocked)"),
        ("2", "search_resources", "Find servers/VMs by cloud, region, type, tags, or keyword", "OpsRamp Resources API (mocked)"),
        ("3", "get_resource_details", "Deep-dive on one resource: config, CPU/mem/disk metrics, tags", "OpsRamp Resources API (mocked)"),
        ("4", "search_incidents", "Find incidents by status, priority, or keyword search", "OpsRamp Incidents API (mocked)"),
        ("5", "investigate_resource", "One-call investigation: resource + alerts + incidents + metrics", "OpsRamp (all APIs combined)"),
        ("6", "get_environment_summary", "Dashboard overview: resource counts, alert/incident breakdown", "OpsRamp (aggregated)"),
    ]

    col_widths_t = [Inches(0.4), Inches(2.3), Inches(4.8), Inches(4.2)]
    table_left = LEFT_MARGIN
    table_top = Inches(1.7)
    row_h = Inches(0.68)

    # Header row
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left, table_top,
        CONTENT_WIDTH, Inches(0.45)
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = HPE_GREEN
    hdr_bg.line.fill.background()

    x_off = table_left + Inches(0.1)
    for j, hdr_text in enumerate(tools_part1[0]):
        hbox = slide.shapes.add_textbox(x_off, table_top + Inches(0.05), col_widths_t[j] - Inches(0.15), Inches(0.35))
        tf = hbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = hdr_text
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE
        r.font.bold = True
        x_off += col_widths_t[j]

    # Data rows
    row_y = table_top + Inches(0.45)
    for i, row in enumerate(tools_part1[1:]):
        bg_color = CARD_BG if i % 2 == 0 else ROW_ALT_BG
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left, row_y,
            CONTENT_WIDTH, row_h
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.color.rgb = CARD_BORDER
        row_bg.line.width = Pt(0.5)

        x_off = table_left + Inches(0.1)
        for j, cell in enumerate(row):
            cbox = slide.shapes.add_textbox(x_off, row_y + Inches(0.06), col_widths_t[j] - Inches(0.15), row_h - Inches(0.1))
            tf = cbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = cell
            if j == 1:  # tool name
                r.font.size = Pt(11)
                r.font.color.rgb = HPE_GREEN
                r.font.bold = True
            elif j == 0:  # number
                r.font.size = Pt(12)
                r.font.color.rgb = WHITE
                r.font.bold = True
            elif j == 3:  # data source
                r.font.size = Pt(10)
                r.font.color.rgb = MEDIUM_GRAY
            else:
                r.font.size = Pt(11)
                r.font.color.rgb = LIGHT_GRAY
            x_off += col_widths_t[j]

        row_y += row_h

    # Bottom note
    note6 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.35), CONTENT_WIDTH, Inches(0.35))
    tf = note6.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "These 6 tools query and retrieve data \u2014 the next 5 tools analyze, correlate, and generate actionable output"
    r.font.size = Pt(12)
    r.font.color.rgb = MEDIUM_GRAY
    r.font.italic = True

    add_hpe_footer(slide, W, H, slide_num=6)


# =============================================================================
# SLIDE 7 — 11 Tools Overview (Part 2: Tools 7–11)
# =============================================================================
def create_slide_7_tools_part2(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 7)
    add_slide_title(slide, "11 AI Tools \u2014 Analysis & Remediation")

    sub7 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.15), Inches(11), Inches(0.4))
    tf = sub7.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "These tools cross-correlate data from OpsRamp + Juniper and produce actionable insights"
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_GRAY

    # Tool table — Part 2
    tools_part2 = [
        ("#", "Tool", "What It Does", "How It Works"),
        ("7", "predict_capacity", "Forecast when CPU/memory/disk hits threshold", "30-day metric history \u2192 linear regression\n\u2192 predicts exhaustion date per resource"),
        ("8", "search_knowledge_base", "Search runbook PDFs for procedures & fixes", "PDF \u2192 chunk (500 chars) \u2192 embed with\nnomic-embed-text \u2192 cosine similarity search"),
        ("9", "correlate_network", "Check if a server issue is caused by the network", "Resolve server IP \u2192 find Juniper switch port\n\u2192 analyze: packet loss, CRC, link flaps, jitter"),
        ("10", "blast_radius", "Map the full impact of an infrastructure failure", "BFS traversal of dependency graph:\nserver \u2192 apps \u2192 downstream \u2192 user groups"),
        ("11", "get_remediation_plan", "Generate step-by-step fix with CLI commands", "Detects issue type \u2192 builds Junos CLI steps\n(clear counters, bounce port, force duplex, etc.)"),
    ]

    col_widths_t2 = [Inches(0.45), Inches(2.5), Inches(4.3), Inches(4.45)]
    table_left = LEFT_MARGIN
    table_top = Inches(1.7)
    row_h = Inches(0.85)

    # Header row
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left, table_top,
        CONTENT_WIDTH, Inches(0.45)
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = HPE_GREEN
    hdr_bg.line.fill.background()

    x_off = table_left + Inches(0.1)
    for j, hdr_text in enumerate(tools_part2[0]):
        hbox = slide.shapes.add_textbox(x_off, table_top + Inches(0.05), col_widths_t2[j] - Inches(0.15), Inches(0.35))
        tf = hbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = hdr_text
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE
        r.font.bold = True
        x_off += col_widths_t2[j]

    # Data rows
    row_y = table_top + Inches(0.45)
    for i, row in enumerate(tools_part2[1:]):
        bg_color = CARD_BG if i % 2 == 0 else ROW_ALT_BG
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left, row_y,
            CONTENT_WIDTH, row_h
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.color.rgb = CARD_BORDER
        row_bg.line.width = Pt(0.5)

        x_off = table_left + Inches(0.1)
        for j, cell in enumerate(row):
            cbox = slide.shapes.add_textbox(x_off, row_y + Inches(0.06), col_widths_t2[j] - Inches(0.15), row_h - Inches(0.1))
            tf = cbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = cell
            if j == 1:  # tool name
                r.font.size = Pt(11)
                r.font.color.rgb = HPE_GREEN
                r.font.bold = True
            elif j == 0:  # number
                r.font.size = Pt(12)
                r.font.color.rgb = WHITE
                r.font.bold = True
            elif j == 3:  # how it works
                r.font.size = Pt(10)
                r.font.color.rgb = MEDIUM_GRAY
            else:
                r.font.size = Pt(11)
                r.font.color.rgb = LIGHT_GRAY
            x_off += col_widths_t2[j]

        row_y += row_h

    # Bottom callout — end-to-end integration
    e2e_card = add_card(slide, LEFT_MARGIN, Inches(5.95), CONTENT_WIDTH, Inches(0.95), border_color=HPE_GREEN)

    e2e_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.98), Inches(11.3), Inches(0.9))
    tf = e2e_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\U0001f517  END-TO-END INTEGRATION \u2014 All 11 tools work as one system"
    r.font.size = Pt(12)
    r.font.color.rgb = HPE_GREEN
    r.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = "The LLM chains tools automatically: OpsRamp alerts \u2192 server metrics \u2192 Juniper switch telemetry \u2192 dependency graph \u2192 PDF runbooks \u2192 CLI fix"
    r2.font.size = Pt(11)
    r2.font.color.rgb = LIGHT_GRAY

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(1)
    r3 = p3.add_run()
    r3.text = "One question triggers up to 6 correlated tool calls across OpsRamp + Juniper + Knowledge Base \u2014 no human orchestration needed"
    r3.font.size = Pt(11)
    r3.font.color.rgb = MEDIUM_GRAY

    add_hpe_footer(slide, W, H, slide_num=7)


# =============================================================================
# SLIDE 8 — Technology Stack
# =============================================================================
def create_slide_8_tech_stack(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 8)
    add_slide_title(slide, "Technology Stack")

    add_section_header(slide, "KEY TECHNOLOGIES", LEFT_MARGIN, Inches(1.5))

    # Technology table data
    tech_rows = [
        ("Go (Golang)", "Core application language \u2014 agent orchestrator, HTTP server, tool execution, MCP server"),
        ("Ollama + llama3.1", "Local LLM inference engine with native tool-calling support for autonomous multi-step reasoning"),
        ("MCP (Model Context Protocol)", "Open standard by Anthropic \u2014 exposes 11 tools to external clients (VS Code Copilot, Claude Desktop) via stdio/HTTP"),
        ("RAG (Retrieval-Augmented Generation)", "PDF runbook ingestion pipeline: text extraction \u2192 chunking (500 chars, 100 overlap) \u2192 embedding \u2192 cosine similarity vector search"),
        ("SSE (Server-Sent Events)", "Real-time streaming of LLM responses and tool progress indicators to the browser \u2014 no WebSockets needed"),
        ("HTML / CSS / JavaScript", "Embedded single-page web UI with ReadableStream SSE parsing and live markdown rendering"),
        ("Docker + Docker Compose", "Containerized deployment \u2014 Ollama + agent in separate containers, supports both native and Docker workflows"),
        ("mcp-go", "Go library (mark3labs/mcp-go) for implementing MCP servers \u2014 converts tool definitions to MCP protocol format"),
    ]

    # Draw table
    table_left = LEFT_MARGIN
    table_top = Inches(2.0)
    col_widths = [Inches(3.2), Inches(8.5)]
    row_height = Inches(0.6)

    # Header row
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left, table_top,
        CONTENT_WIDTH, Inches(0.45)
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = HPE_GREEN
    hdr_bg.line.fill.background()

    hdr_tech = slide.shapes.add_textbox(table_left + Inches(0.15), table_top + Inches(0.05), col_widths[0], Inches(0.35))
    tf = hdr_tech.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Technology"
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True

    hdr_use = slide.shapes.add_textbox(table_left + col_widths[0] + Inches(0.15), table_top + Inches(0.05), col_widths[1], Inches(0.35))
    tf = hdr_use.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Usage in HPE AI Autopilot"
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Data rows
    row_y = table_top + Inches(0.45)
    for i, (tech, usage) in enumerate(tech_rows):
        bg_color = CARD_BG if i % 2 == 0 else ROW_ALT_BG

        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            table_left, row_y,
            CONTENT_WIDTH, row_height
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.color.rgb = CARD_BORDER
        row_bg.line.width = Pt(0.5)

        tech_box = slide.shapes.add_textbox(table_left + Inches(0.15), row_y + Inches(0.08), col_widths[0] - Inches(0.15), row_height - Inches(0.1))
        tf = tech_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tech
        r.font.size = Pt(13)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True

        use_box = slide.shapes.add_textbox(table_left + col_widths[0] + Inches(0.15), row_y + Inches(0.08), col_widths[1] - Inches(0.3), row_height - Inches(0.1))
        tf = use_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = usage
        r.font.size = Pt(12)
        r.font.color.rgb = LIGHT_GRAY

        row_y += row_height

    add_hpe_footer(slide, W, H, slide_num=8)


# =============================================================================
# SLIDE 9 — What's Next / Roadmap (NEW)
# =============================================================================
def create_slide_9_roadmap(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 9)
    add_slide_title(slide, "What\u2019s Next \u2014 Roadmap")

    # Subtitle
    sub = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.15), Inches(11), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "From innovation demo to strategic initiative \u2014 the path to production"
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_GRAY

    # --- Left column: Near-Term ---
    add_section_header(slide, "NEAR-TERM (Q2 2026)", LEFT_MARGIN, Inches(1.8))

    near_items = [
        ("Real OpsRamp API Integration", "Replace mocked data layer with live OpsRamp V3 API calls \u2014\nauthenticate via OAuth, query real alerts/resources/incidents"),
        ("Real Juniper Mist API", "Connect to production Juniper Mist cloud for live switch stats,\nport telemetry, and network event correlation"),
        ("Expanded Runbook Library", "Ingest additional PDF/Markdown runbooks covering server,\nnetwork, storage, and application troubleshooting"),
        ("LLM Model Upgrades", "Evaluate larger models (llama3.1 70B, Mixtral) for improved\ntool-calling accuracy and reasoning depth"),
    ]

    item_y = Inches(2.2)
    for title, desc in near_items:
        add_card(slide, LEFT_MARGIN, item_y, Inches(5.5), Inches(0.9))

        tbox = slide.shapes.add_textbox(Inches(1.0), item_y + Inches(0.05), Inches(5.1), Inches(0.85))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"\u2713  {title}"
        r.font.size = Pt(13)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = f"    {desc}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = LIGHT_GRAY

        item_y += Inches(0.95)

    # Divider
    add_vertical_divider(slide, Inches(6.6), Inches(1.8))

    # --- Right column: Medium-Term ---
    add_section_header(slide, "MEDIUM-TERM (H2 2026)", RIGHT_COL, Inches(1.8))

    medium_items = [
        ("Multi-Tenant Support", "Tenant-aware API calls and data isolation \u2014 serve multiple\ncustomer environments from a single agent deployment"),
        ("Production Deployment", "Containerized deployment on HPE GreenLake infrastructure\nwith CI/CD pipeline, monitoring, and auto-scaling"),
        ("Additional Data Sources", "Integrate Aruba Central, iLO/iSM, and storage array APIs\nfor full-stack infrastructure correlation"),
        ("Self-Healing Actions", "Graduate from \u201cguided remediation\u201d to \u201capproved auto-fix\u201d \u2014\nexecute safe remediation steps with human approval gates"),
    ]

    item_y = Inches(2.2)
    for title, desc in medium_items:
        add_card(slide, RIGHT_COL, item_y, Inches(5.5), Inches(0.9))

        tbox = slide.shapes.add_textbox(Inches(7.4), item_y + Inches(0.05), Inches(5.1), Inches(0.85))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"\u25b6  {title}"
        r.font.size = Pt(13)
        r.font.color.rgb = HPE_GREEN
        r.font.bold = True
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = f"    {desc}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = LIGHT_GRAY

        item_y += Inches(0.95)

    # Bottom vision callout
    callout = add_card(slide, LEFT_MARGIN, Inches(6.2), CONTENT_WIDTH, Inches(0.7), border_color=HPE_GREEN)

    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\U0001f680  VISION: Autonomous IT Operations \u2014 AI that detects, investigates, correlates, and resolves infrastructure issues end-to-end"
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.bold = True

    add_hpe_footer(slide, W, H, slide_num=9)


# =============================================================================
# SLIDE 10 — UI Preview / Screenshots (NEW)
# =============================================================================
def create_slide_10_screenshots(prs, W, H):
    slide = add_slide_boilerplate(prs, W, H, 10)
    add_slide_title(slide, "UI Preview \u2014 AI Autopilot in Action")

    sub = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.15), Inches(11), Inches(0.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Web UI conversation flow and MCP server integration with VS Code / Claude Desktop"
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_GRAY

    # --- Left: Web UI screenshot ---
    add_section_header(slide, "WEB UI", LEFT_MARGIN, Inches(1.7))

    # Screenshot frame
    web_frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(2.1), Inches(5.8), Inches(4.0)
    )
    web_frame.fill.solid()
    web_frame.fill.fore_color.rgb = CARD_BG
    web_frame.line.color.rgb = HPE_GREEN
    web_frame.line.width = Pt(1.5)

    # Insert web-ui.png
    web_ui_path = os.path.join(SCREENSHOTS_DIR, "web-ui.png")
    if os.path.exists(web_ui_path):
        slide.shapes.add_picture(
            web_ui_path,
            LEFT_MARGIN + Inches(0.1), Inches(2.2),
            Inches(5.6), Inches(3.8)
        )
    else:
        # Fallback text if image not found
        fb = slide.shapes.add_textbox(LEFT_MARGIN + Inches(1.0), Inches(3.5), Inches(3.8), Inches(1.0))
        tf = fb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "[web-ui.png not found]\nPlace screenshot in screenshots/ folder"
        r.font.size = Pt(14)
        r.font.color.rgb = MEDIUM_GRAY

    # Caption
    web_cap = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.15), Inches(5.8), Inches(0.4))
    tf = web_cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Streaming conversation with tool call indicators and live markdown"
    r.font.size = Pt(10)
    r.font.color.rgb = MEDIUM_GRAY

    # --- Right: MCP Server screenshot ---
    add_section_header(slide, "MCP SERVER", RIGHT_COL, Inches(1.7))

    mcp_frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.0)
    )
    mcp_frame.fill.solid()
    mcp_frame.fill.fore_color.rgb = CARD_BG
    mcp_frame.line.color.rgb = HPE_GREEN
    mcp_frame.line.width = Pt(1.5)

    # Insert mcp-server.png
    mcp_path = os.path.join(SCREENSHOTS_DIR, "mcp-server.png")
    if os.path.exists(mcp_path):
        slide.shapes.add_picture(
            mcp_path,
            Inches(7.1), Inches(2.2),
            Inches(5.6), Inches(3.8)
        )
    else:
        fb = slide.shapes.add_textbox(Inches(8.0), Inches(3.5), Inches(3.8), Inches(1.0))
        tf = fb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "[mcp-server.png not found]\nPlace screenshot in screenshots/ folder"
        r.font.size = Pt(14)
        r.font.color.rgb = MEDIUM_GRAY

    # Caption
    mcp_cap = slide.shapes.add_textbox(Inches(7.0), Inches(6.15), Inches(5.8), Inches(0.4))
    tf = mcp_cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "MCP protocol exposes all 11 tools to VS Code Copilot and Claude Desktop"
    r.font.size = Pt(10)
    r.font.color.rgb = MEDIUM_GRAY

    add_hpe_footer(slide, W, H, slide_num=10)


# =============================================================================
# SLIDE 11 — Demo & Thank You
# =============================================================================
def create_slide_11_demo_thankyou(prs, W, H):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_hpe_background(slide, HPE_DARKER)
    add_hpe_header_bar(slide, W)
    add_green_accent_shape(slide, Inches(0), Inches(0.06), ACCENT_WIDTH, H - Inches(0.56))

    # Large centered "Live Demo" text
    demo_box = slide.shapes.add_textbox(Inches(0), Inches(1.5), W, Inches(1.2))
    tf = demo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Live Demo"
    r.font.size = Pt(56)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Decorative separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(2.9),
        Inches(2.3), Inches(0.04)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = HPE_GREEN
    sep.line.fill.background()

    # Demo label — more presentable than raw localhost
    label_box = slide.shapes.add_textbox(Inches(0), Inches(3.1), W, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "HPE AI Autopilot \u2014 Local Demo"
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Demo URL (secondary)
    url_box = slide.shapes.add_textbox(Inches(0), Inches(3.6), W, Inches(0.4))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "http://localhost:8080"
    r.font.size = Pt(14)
    r.font.color.rgb = HPE_GREEN

    # Thank you
    ty_box = slide.shapes.add_textbox(Inches(0), Inches(4.5), W, Inches(1.0))
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank You"
    r.font.size = Pt(44)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Presenter
    pres_box = slide.shapes.add_textbox(Inches(0), Inches(5.5), W, Inches(0.5))
    tf = pres_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Vishnu Dutt Kidambi"
    r.font.size = Pt(18)
    r.font.color.rgb = LIGHT_GRAY

    # Q&A note
    qa_box = slide.shapes.add_textbox(Inches(0), Inches(6.1), W, Inches(0.4))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Questions & Discussion"
    r.font.size = Pt(16)
    r.font.color.rgb = MEDIUM_GRAY

    add_hpe_footer(slide, W, H, slide_num=11, show_confidential=False)


# =============================================================================
# Main
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # Create all 11 slides
    create_slide_1_title(prs, W, H)
    create_slide_2_problem_solution(prs, W, H)
    create_slide_3_architecture(prs, W, H)
    create_slide_4_data_connection(prs, W, H)
    create_slide_5_e2e_flow(prs, W, H)
    create_slide_6_tools_part1(prs, W, H)
    create_slide_7_tools_part2(prs, W, H)
    create_slide_8_tech_stack(prs, W, H)
    create_slide_9_roadmap(prs, W, H)
    create_slide_10_screenshots(prs, W, H)
    create_slide_11_demo_thankyou(prs, W, H)

    # Save
    output_path = "HPE_AI_Autopilot_Demo.pptx"
    prs.save(output_path)
    print(f"\u2705 Presentation saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    print()
    print("\u2550" * 65)
    print("  SLIDE OVERVIEW")
    print("\u2550" * 65)
    print("   1. Title \u2014 HPE AI Autopilot (Vishnu Dutt Kidambi)")
    print("   2. Why AI Autopilot? (Challenge vs Solution)")
    print("   3. Architecture & Data Sources")
    print("   4. OpsRamp \u2194 Juniper: IP Join Key, Port Mapping, Juniper Scope")
    print('   5. End-to-End Flow: "Why is GreenLake Portal Slow?" (6-step walkthrough)')
    print("   6. 11 AI Tools \u2014 Data & Investigation (Tools 1\u20136)")
    print("   7. 11 AI Tools \u2014 Analysis & Remediation (Tools 7\u201311)")
    print("   8. Technology Stack (MCP, RAG, Go, SSE, Docker, etc.)")
    print("   9. What\u2019s Next \u2014 Roadmap (Near-Term & Medium-Term)")
    print("  10. UI Preview \u2014 Screenshots (Web UI + MCP Server)")
    print("  11. Live Demo & Thank You")
    print("\u2550" * 65)


if __name__ == "__main__":
    main()
